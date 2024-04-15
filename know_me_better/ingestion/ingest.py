import pdb
import os
from github import Github
from langchain.docstore.document import Document
from langchain_community.vectorstores import SupabaseVectorStore

from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import CharacterTextSplitter
from supabase.client import Client, create_client

from dotenv import load_dotenv

import requests
import fitz  # PyMuPDF library for PDF files
import pptx  # python-pptx library for PowerPoint files


def get_pdf_content_from_github(repo_owner, repo_name, file_path):
    """
    Retrieves the content of a PDF file from a GitHub repository.

    Args:
        repo_owner (str): The owner of the GitHub repository.
        repo_name (str): The name of the GitHub repository.
        file_path (str): The path to the PDF file within the repository.

    Returns:
        str: The content of the PDF file.
    """
    url = f"https://raw.githubusercontent.com/{repo_owner}/{repo_name}/master/{file_path}"
    response = requests.get(url)

    if response.status_code == 200:
        content = ""
        if file_path.endswith(".pdf"):
            with fitz.open(stream=response.content, filetype="pdf") as pdf_file:
                for page in pdf_file:
                    content += page.get_text()
            return content

        elif file_path.endswith(".pptx"):
            with pptx.Presentation.open(stream=response.content) as presentation:
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text
            return content
        else:
            print(f"File type not supported: {file_path}")
            return response.text
    else:
        raise Exception(f"Failed to retrieve file from GitHub: {response.status_code} - {response.text}")


def main():
    # Load environment variables
    load_dotenv()

    # Set up GitHub API
    GITHUB_TOKEN = os.getenv("GITHUB_TOKEN")
    github_user_name = "jagan-shanmugam"
    github_repo_name = "Course-Work"
    embedding_model = "all-mpnet-base-v2"
    chunk_size = 1024
    chunk_overlap = 256
    
    text_splitter = CharacterTextSplitter(chunk_size=chunk_size,
                                           chunk_overlap=chunk_overlap)
    g = Github(GITHUB_TOKEN)

    supabase_url = os.environ.get("SUPABASE_URL")
    supabase_key = os.environ.get("SUPABASE_SERVICE_KEY")
    supabase: Client = create_client(supabase_url, supabase_key)


    # Set up OpenAI Embeddings
    # embeddings = OpenAIEmbeddings()
    # model = SentenceTransformer(embedding_model)

    embeddings = HuggingFaceEmbeddings(model_name=embedding_model)

    # Set up vector store
    # vector_store = SupabaseVectorStore(SUPABASE_URL, SUPABASE_KEY, table_name="github_data")
    # vector_store = SupabaseVectorStore(
    #     embedding=embeddings,
    #     client=supabase,
    #     table_name="documents",
    #     query_name="match_documents",
    # )

    # Specify the GitHub repository to read from
    repo = g.get_repo(f"{github_user_name}/{github_repo_name}")
    #
    # Iterate through files in the repository
    contents = repo.get_contents("")
    docs = []
    while contents:
        file = contents.pop(0)
        if file.type == "dir":
            contents.extend(repo.get_contents(file.path))
        else:
            
            try:
                if file.path.endswith(".md") or file.path.endswith(".txt"):
                    content = repo.get_contents(file.path).decoded_content.decode("utf-8")
                    
                    texts = text_splitter.split_text(content)
                    docs.extend([Document(page_content=t) for t in texts])
                    print(f"Processed file: {file.path}")

                elif file.path.endswith(".pdf"):
                    content = get_pdf_content_from_github(github_user_name, github_repo_name, file.path)
                    # with fitz.open(file.path) as pdf_file:
                    #     content = ""
                    #     for page in pdf_file:
                    #         content += page.get_text()

                    texts = text_splitter.split_text(content)
                    docs.extend([Document(page_content=t) for t in texts])
                    print(f"Processed file: {file.path}")

                elif file.path.endswith(".pptx"):
                    # presentation = pptx.Presentation(file.path)
                    content = get_pdf_content_from_github(github_user_name, github_repo_name, file.path)

                    # content = ""
                    # for slide in presentation.slides:
                    #     for shape in slide.shapes:
                    #         if hasattr(shape, "text"):
                    #             content += shape.text
                    
                    texts = text_splitter.split_text(content)
                    docs.extend([Document(page_content=t) for t in texts])
                    
                    print(f"Processed file: {file.path}")

            except Exception as e:
                print(f"Error processing file {file.path}: {e}")

    # Save the vector store to Supabase
    # vector_store.persist()
    pdb.set_trace()
    vector_store = SupabaseVectorStore.from_documents(
        docs,
        embeddings,
        client=supabase,
        table_name="documents",
        query_name="match_documents",
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap
    )
    print("Vector store saved to Supabase.")

    pdb.set_trace()

if __name__ == "__main__":
    main()
